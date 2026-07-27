"""PPTX parser (lazy). Uses python-pptx when available; emits slides
as theory blocks. No MCQ layout is expected in PPT so the slides are
treated as training material that gets AI-classified and indexed.
"""
from __future__ import annotations

import io
import time
from pathlib import Path
from typing import List

from .dataclasses import ParsedDocument, ParsedTheory
from .text_utils import clean_text


class PPTXParser:
    def parse(self, path: str) -> ParsedDocument:
        start = time.time()
        fn = Path(path).name
        doc = ParsedDocument(
            source_filename=fn,
            file_format="pptx",
            detected_type="theory",
            parser_used="pptx",
        )
        try:
            from pptx import Presentation  # type: ignore
        except Exception:
            doc.errors.append("python-pptx is not installed.")
            return doc
        try:
            presentation = Presentation(path)
        except Exception as exc:
            doc.errors.append(f"Failed to open PPTX: {exc}")
            return doc

        slides: List[ParsedTheory] = []
        for idx, slide in enumerate(presentation.slides):
            buf: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    buf.append(shape.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text:
                                buf.append(cell.text)
            text = clean_text("\n".join(buf))
            slides.append(
                ParsedTheory(
                    position_index=idx,
                    heading=f"Slide {idx + 1}",
                    body_text=text,
                    block_type="mixed",
                    raw_text=text,
                )
            )
        doc.theory_blocks = slides
        doc.duration_ms = int((time.time() - start) * 1000)
        return doc
